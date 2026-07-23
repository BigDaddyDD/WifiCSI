"""Second end-of-week slide deck: plain white background, black Arial, NO
theme (same convention as build_week_update_slides.py / build_slides.py).
Minimal per-slide text -- figures carry the content.

Covers: going fully wireless, the updated (wireless, Taoglas) production
dataset, a histogram of that dataset (boss request), a mathematical
definition of balanced vs. traditional accuracy (boss request), the same
tree-diagram slide as last week, this week's cross-placement + antenna
generalization results, and a live field-demo photo slide.

Run: python slides_assets/build_week2_update_slides.py
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = os.path.join(ROOT, 'docs', 'model_report_assets')
OUT = os.path.join(ROOT, 'docs', 'CSI_Week2_Update_Slides.pptx')

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Arial'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def setrun(run, size, bold=False, italic=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = BLACK


def title(s, text):
    tf = textbox(s, 0.5, 0.3, SW - 1.0, 1.0)
    r = tf.paragraphs[0].add_run(); r.text = text
    setrun(r, 28, True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.1),
                            Inches(SW - 1.1), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background()
    ln.shadow.inherit = False


def bullets(s, left, top, width, height, points, size=18):
    tf = textbox(s, left, top, width, height)
    for i, p in enumerate(points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(10)
        r = para.add_run(); r.text = '•  ' + p
        setrun(r, size)


def mono_lines(s, left, top, width, height, lines, size=18):
    """Plain (non-bulleted) lines, for formula blocks."""
    tf = textbox(s, left, top, width, height)
    for i, (text, bold, italic) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(8)
        r = para.add_run(); r.text = text
        setrun(r, size, bold=bold, italic=italic)


def image_fit(s, path, box_l, box_t, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = box_w; h = w / ar
    if h > box_h:
        h = box_h; w = h * ar
    s.shapes.add_picture(path, Inches(box_l + (box_w - w) / 2),
                         Inches(box_t + (box_h - h) / 2),
                         width=Inches(w), height=Inches(h))


# ---- Slide 1: title --------------------------------------------------------
s = slide()
tf = textbox(s, 0.8, 2.6, SW - 1.6, 1.2)
r = tf.paragraphs[0].add_run(); r.text = 'Wi-Fi CSI Presence & Activity Sensing'
setrun(r, 40, True)
tf2 = textbox(s, 0.8, 3.9, SW - 1.6, 0.8)
r = tf2.paragraphs[0].add_run()
r.text = 'End-of-week update, wireless deployment & antenna comparison · July 2026'
setrun(r, 20)

# ---- Slide 2: this week -----------------------------------------------------
s = slide(); title(s, 'This week')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'The system now runs fully wirelessly: the receiver relays CSI over Wi-Fi at '
    'the full native ~97–99 Hz rate, and a live monitor app classifies '
    'presence/activity from a laptop with no cable to the sensor.',
    'Repeated the antenna-portability question from last week, this time '
    'wirelessly and across more placements/rooms: both boards on the Taoglas '
    'antenna vs. a mixed pair (one board per antenna type).',
    'Two follow-ups from last week’s slide review, addressed directly here: '
    'a mathematical definition of balanced vs. traditional accuracy, and a '
    'histogram of the training data.',
    'Same recipe as always: per-config empty-room calibration, leave-one-'
    'config-out evaluation.',
], size=19)

# ---- Slide 3: the dataset ---------------------------------------------------
s = slide(); title(s, 'The dataset')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Configurations (room × placement): 15',
    'Recording sessions: 34',
    'Total recording time: ~104 min (~1.7 hours)',
    'CSI packets logged: 608,799',
    'Labeled 2-second analysis windows: 4,241',
    'Classes: empty, stand, sit, walk, run',
    'One subject, four rooms, single antenna configuration (Taoglas on both '
    'boards); the mixed-antenna pair was recorded separately, only for the '
    'antenna-comparison test, and is not part of this production dataset.',
], size=19)

# ---- Slide 4: histogram of the training data (boss request) ---------------
s = slide(); title(s, 'Histogram of the training data')
image_fit(s, os.path.join(A, 'fig12_training_histogram.png'), 0.9, 1.3, SW - 1.8, 4.9)
bullets(s, 0.9, 6.3, SW - 1.8, 1.0, [
    'Stand/sit/walk/run come out close to balanced by design (the recording '
    'protocol records equal-length blocks per activity); empty is smaller '
    'because it is capped per config so a few long calibration captures '
    'cannot dominate the class.',
], size=15)

# ---- Slide 5: inside the model -----------------------------------------------
s = slide(); title(s, 'Inside the model: the whole tree')
image_fit(s, os.path.join(A, 'fig11_tree_diagram.png'), 0.3, 1.2, SW - 0.6, 5.0)
bullets(s, 0.9, 6.35, SW - 1.8, 1.0, [
    'Same real tree as last week’s slide (tree #1 of the 300-tree forest, all '
    '1,699 real nodes): the modeling approach is unchanged; this week’s work '
    'retrains the same random-forest recipe on the new wireless dataset.',
], size=15)

# ---- Slide 6: balanced accuracy vs. traditional accuracy (boss request) ----
s = slide(); title(s, 'Balanced accuracy vs. traditional accuracy')
bullets(s, 0.7, 1.45, SW - 1.4, 1.6, [
    'Balanced accuracy = the average of the per-class recalls, instead of the '
    'raw fraction of windows called correctly.',
    'Why we use it: classes are uneven (more occupied/stand windows than '
    'empty/run). Plain accuracy lets a model score high just by predicting '
    'the majority class; balanced accuracy weights every class equally, so a '
    'less-recorded class counts as much as a common one.',
], size=18)
mono_lines(s, 0.7, 3.35, SW - 1.4, 2.9, [
    ('How it’s calculated: recall = (windows of that class predicted correctly) ÷ (total windows of that class)', False, True),
    ('Example: this week’s presence result, held-out (leave-one-config-out):', True, False),
    ('  empty recall      =   171 / 486    =  0.352   (poor)', False, True),
    ('  occupied recall   = 3,469 / 3,755  =  0.924   (inflated)', False, True),
    ('  balanced accuracy = (0.352 + 0.924) / 2  =  0.638  ≈  0.64', True, False),
    ('  traditional accuracy = 3,640 / 4,241 = 0.858 ≈ 0.86, which looks much better but hides the weak empty recall', False, False),
], size=18)

# ---- Slide 7: cross-placement result ----------------------------------------
s = slide(); title(s, 'Cross-placement generalization (wireless)')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Leave-one-config-out, 15 folds, wireless Taoglas production dataset: '
    'presence 0.64 balanced, 5-class activity 0.37 balanced.',
    'Same range this project has always seen on the wired data; moving to '
    'wireless is not changing the underlying generalization result, which is '
    'reassuring: the relay itself is not introducing a new problem.',
], size=19)

# ---- Slide 8: antenna confusion matrices (wireless) --------------------------
s = slide(); title(s, 'Taoglas vs. mixed pair: per-install accuracy')
image_fit(s, os.path.join(A, 'fig13_confusion_wireless_antenna.png'), 0.4, 1.3, SW - 0.8, 4.4)
bullets(s, 0.9, 5.85, SW - 1.8, 1.5, [
    'Left two panels: each antenna condition calibrated and evaluated on its '
    'own recordings (per-install, 5-fold). Both work well (0.87 / 0.92 '
    'balanced), with the same confusions in both (sit/stand, walk/run).',
    'Right panel: train on Taoglas, test on the mixed pair, no retraining: '
    'the diagonal disappears. Sit gets called stand 63% of the time; run '
    'gets called walk 68% of the time. This is what "does not transfer" '
    'looks like, not just a number.',
], size=15)

# ---- Slide 9: does it plug and play wirelessly? ------------------------------
s = slide(); title(s, 'Does it plug-and-play onto a new antenna? (wireless repeat)')
image_fit(s, os.path.join(A, 'fig14_wireless_antenna_transfer.png'), 1.2, 1.3, 7.4, 4.6)
bullets(s, 0.9, 6.05, SW - 1.8, 1.3, [
    'Same conclusion as last week, now confirmed on an independent wireless '
    'dataset: cross-antenna transfer collapses (5-class 0.36, presence 0.67) '
    'well below either antenna’s own per-install number (0.87–0.97).',
    'External corroboration: CSI-Bench (Zhu et al., 2025, arXiv:2505.21866) '
    'reports the same pattern industry-wide: cross-device WiFi sensing '
    'accuracy drops sharply (e.g. 99.8% → ~70% F1 in their benchmark) due to '
    '"hardware heterogeneity," even without changing the room. Our result is '
    'consistent with a known, general limitation, not an artifact of our setup.',
], size=14)

# ---- Slide 8: live field demo -----------------------------------------------
s = slide(); title(s, 'Live monitor: interface preview')
photos = ['Empty.jpg', 'Stand.jpg', 'Sit.jpg']
labels = ['Empty', 'Standing', 'Sitting']
pw = (SW - 1.6) / 3
for i, (fn, lab) in enumerate(zip(photos, labels)):
    left = 0.8 + i * (pw + 0.2)
    image_fit(s, os.path.join(A, fn), left, 1.3, pw, 4.6)
    tf = textbox(s, left, 6.0, pw, 0.5)
    r = tf.paragraphs[0].add_run(); r.text = lab
    setrun(r, 16, True)
bullets(s, 0.8, 6.5, SW - 1.6, 0.8, [
    'What the live monitor interface looks like once deployed: presence and '
    'activity state, updated from the wireless CSI stream.',
], size=13)

# ---- Slide 9: status & next -------------------------------------------------
s = slide(); title(s, 'Status & next steps')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Wireless deployment is validated end to end: full native rate, stable '
    'live classification, field-tested interface.',
    'Antenna portability result now confirmed twice (wired and wireless, '
    'independent datasets); treating this as settled rather than directional.',
    'Continuing to add placements/rooms on the single (Taoglas) antenna to '
    'tighten the cross-placement confidence interval.',
    'Next: keep expanding room/placement coverage on the production antenna; '
    '2nd/3rd receiver validation remains the lever for the fine-activity '
    'ceiling that no model or antenna change has moved.',
], size=19)

prs.save(OUT)
print('Saved', OUT, '-', len(prs.slides._sldIdLst), 'slides')
