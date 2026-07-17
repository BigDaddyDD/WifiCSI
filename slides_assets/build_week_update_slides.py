"""Brief end-of-week slide deck: plain white background, black Arial, NO theme
(same convention as build_slides.py). Minimal per-slide text — figures carry
the content; text is a short placeholder the user rewrites in their own voice.

Covers: this week's stock-vs-Taoglas antenna portability test, why the
project still uses a random forest (fig7 model comparison), a look inside the
model (tree diagram + feature importance), and the antenna results.

Run: python slides_assets/build_week_update_slides.py
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = os.path.join(ROOT, 'docs', 'model_report_assets')
OUT = os.path.join(ROOT, 'docs', 'CSI_Week_Update_Slides.pptx')

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Arial'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def setrun(run, size, bold=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = BLACK


def title(s, text):
    tf = textbox(s, 0.5, 0.3, SW - 1.0, 1.0)
    r = tf.paragraphs[0].add_run(); r.text = text
    setrun(r, 28, True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.1),
                            Inches(SW - 1.1), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background()
    ln.shadow.inherit = False


def bullets(s, left, top, width, height, points, size=18):
    tf = textbox(s, left, top, width, height)
    for i, p in enumerate(points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(10)
        r = para.add_run(); r.text = '•  ' + p
        setrun(r, size)


def image_fit(s, path, box_l, box_t, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = box_w; h = w / ar
    if h > box_h:
        h = box_h; w = h * ar
    s.shapes.add_picture(path, Inches(box_l + (box_w - w) / 2),
                         Inches(box_t + (box_h - h) / 2),
                         width=Inches(w), height=Inches(h))


# ---- Slide 1: title --------------------------------------------------------
s = slide()
tf = textbox(s, 0.8, 2.6, SW - 1.6, 1.2)
r = tf.paragraphs[0].add_run(); r.text = 'Wi-Fi CSI Presence & Activity Sensing'
setrun(r, 40, True)
tf2 = textbox(s, 0.8, 3.9, SW - 1.6, 0.8)
r = tf2.paragraphs[0].add_run()
r.text = 'End-of-week update — model comparison & antenna portability · July 2026'
setrun(r, 20)

# ---- Slide 2: this week -----------------------------------------------------
s = slide(); title(s, 'This week')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Question from leadership: does the trained model plug-and-play onto a '
    'different antenna, or a different capture device entirely?',
    'Ran the same placement/activity protocol with two antennas — the stock '
    'ESP32-C3 antenna and a Taoglas unit — at three positions in the same room.',
    'Also ran a fixed-protocol comparison across 10 model families (classical '
    'and deep) to answer a standing question: is random forest still the right '
    'choice, or would a fancier model do better?',
    'Both experiments below use the same recipe as always: per-config empty-'
    'room calibration, leave-one-config-out evaluation.',
], size=19)

# ---- Slide 3: the dataset ---------------------------------------------------
s = slide(); title(s, 'The dataset')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Configurations (placement × orientation): 14',
    'Recording sessions: 20',
    'Total recording time: ~113 min (~1.9 hours)',
    'CSI packets logged: 669,353',
    'Labeled 2-second analysis windows: 4,971',
    'Classes: empty, stand, sit, walk, run',
    'One subject, two library study rooms — a weak-link room and one '
    'range-limited configuration were excluded (documented separately), so '
    'every number above is from links that held up.',
], size=19)

# ---- Slide 4: model comparison -----------------------------------------------
s = slide(); title(s, 'Why we still use a random forest')
image_fit(s, os.path.join(A, 'fig7_model_comparison.png'), 0.6, 1.3, SW - 1.2, 4.3)
bullets(s, 0.9, 5.7, SW - 1.8, 1.6, [
    'Ten model families, same data / calibration / evaluation — only the '
    'classifier changes.',
    'Cross-placement activity is model-independent (0.37–0.46 balanced across '
    'every model); no algorithm breaks out.',
    'Presence favors the engineered features (0.75–0.84) over raw deep nets '
    '(1D-CNN 0.59, GRU 0.64) — random forest stays at or near the top for both.',
], size=15)

# ---- Slide 5: inside one real tree ------------------------------------------
s = slide(); title(s, 'Inside the model: the whole tree')
image_fit(s, os.path.join(A, 'fig11_tree_diagram.png'), 0.3, 1.2, SW - 0.6, 5.0)
bullets(s, 0.9, 6.35, SW - 1.8, 1.0, [
    'This is tree #1 of the actual 300-tree production model — all 1,699 real '
    'nodes, nothing cut off; the top rows show real thresholds, the fan below '
    'is the same real tree continuing down to depth 36.',
], size=15)

# ---- Slide 6: feature importance --------------------------------------------
s = slide(); title(s, 'Which features the model actually uses')
image_fit(s, os.path.join(A, 'fig10_feature_importance.png'), 0.9, 1.4, 7.6, 5.3)
bullets(s, 8.7, 2.3, 4.1, 4.5, [
    'Top features are all deviation-from-empty measures (max/mean relative '
    'amplitude change, overall L2 deviation) — confirms the empty-baseline '
    'calibration is doing the heavy lifting.',
    'Individual subcarriers matter too, but no single one dominates — the '
    'signal is spread across the channel, which is why the model needs many '
    'subcarriers rather than one.',
])

# ---- Slide 7: antenna confusion matrices ------------------------------------
s = slide(); title(s, 'Stock vs. Taoglas: per-install accuracy')
image_fit(s, os.path.join(A, 'fig8_confusion_antenna.png'), 0.5, 1.3, SW - 1.0, 4.6)
bullets(s, 0.9, 6.0, SW - 1.8, 1.3, [
    'Each antenna calibrated and evaluated on its own recordings (per-install, '
    '5-fold): stock 0.72 balanced, Taoglas 0.77 balanced — both antennas work '
    'fine on their own.',
    'Taoglas link was consistently stronger (~10 dB better RSSI) and scored '
    'slightly higher; same confusions in both (sit/stand, walk/run).',
])

# ---- Slide 8: does it plug and play? ----------------------------------------
s = slide(); title(s, 'Does it plug-and-play onto a new antenna?')
image_fit(s, os.path.join(A, 'fig9_antenna_transfer.png'), 1.2, 1.3, 7.4, 5.0)
bullets(s, 0.9, 6.4, SW - 1.8, 1.0, [
    'No. The ESP32 stock antenna and the Taoglas antenna each work well '
    'calibrated on their own data — but a model trained on one and dropped '
    'onto the other, with no retraining, collapses to near-random.',
    'A new antenna or device is not a drop-in swap: CSI is hardware-specific. '
    'What transfers is the method, not the trained model.',
], size=15)

# ---- Slide 9: status & next -------------------------------------------------
s = slide(); title(s, 'Status & next steps')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Model choice is settled: random forest remains the right tool at this data '
    'scale — the ceiling is the single-antenna sensor, not the algorithm.',
    'Antenna portability question answered directly: retraining on the new '
    'antenna’s own data works, reusing the old antenna’s trained model does '
    'not — set that expectation with any new hardware.',
    'Continuing antenna recordings at more positions to firm up the result.',
    'Next: 2nd/3rd receiver validation (spare hardware on hand) — the lever for '
    'both coverage and fine-activity accuracy.',
], size=19)

prs.save(OUT)
print('Saved', OUT, '-', len(prs.slides._sldIdLst), 'slides')
