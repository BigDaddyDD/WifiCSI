#!/usr/bin/env python3
"""Brief Phase-A slide deck: plain white background, black text, NO theme.
Each content slide = one existing (color) figure + a few short talking points.
Run report_figs.py first so docs/report_assets/*.png exist."""
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

A = os.path.join('docs', 'report_assets')
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Arial'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    s = prs.slides.add_slide(BLANK)
    # explicit white background rectangle behind everything (no theme reliance)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def setrun(run, size, bold=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = BLACK


def title(s, text):
    tf = textbox(s, 0.5, 0.3, SW - 1.0, 1.0)
    r = tf.paragraphs[0].add_run(); r.text = text
    setrun(r, 30, True)
    # thin black underline rule
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.15),
                            Inches(SW - 1.1), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background()
    ln.shadow.inherit = False


def bullets(s, left, top, width, height, points, size=18):
    tf = textbox(s, left, top, width, height)
    for i, p in enumerate(points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(10)
        r = para.add_run(); r.text = '•  ' + p
        setrun(r, size)


def image_fit(s, path, box_l, box_t, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = box_w; h = w / ar
    if h > box_h:
        h = box_h; w = h * ar
    s.shapes.add_picture(path, Inches(box_l + (box_w - w) / 2),
                         Inches(box_t + (box_h - h) / 2),
                         width=Inches(w), height=Inches(h))


# ---- Slide 1: title --------------------------------------------------------
s = slide()
tf = textbox(s, 0.8, 2.6, SW - 1.6, 1.2)
r = tf.paragraphs[0].add_run(); r.text = 'Wi-Fi CSI Presence & Activity Sensing'
setrun(r, 40, True)
tf2 = textbox(s, 0.8, 3.9, SW - 1.6, 0.8)
r = tf2.paragraphs[0].add_run()
r.text = 'Phase A interim — single-link characterization · July 2026'
setrun(r, 20)

# ---- Slide 2: setup (text only) -------------------------------------------
s = slide(); title(s, 'Setup')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'One Wi-Fi link: two ESP32-C3 boards (2.4 GHz, single antenna). One transmits '
    'a steady ~100 Hz stream; the other reports CSI for 52 subcarriers at ~97 Hz.',
    'Library study room. Five configurations = different board placements and '
    'antenna orientations (the system must not depend on one exact placement).',
    'Each configuration runs the same script: empty / stand / sit / walk / run.  '
    'Five sessions, ~1,400 analysis windows (2 s each).',
    'Method: per-window features → random forest, calibrated against the room’s '
    'own empty baseline. Tested on a configuration never seen in training '
    '(leave-one-session-out) — the honest measure of generalization.',
], size=19)

# ---- Slide 3: methods figure ----------------------------------------------
s = slide(); title(s, 'Presence works; raw CSI alone does not')
image_fit(s, os.path.join(A, 'fig1_methods.png'), 0.6, 1.35, SW - 1.2, 4.05)
bullets(s, 0.9, 5.6, SW - 1.8, 1.7, [
    'Presence on an unseen configuration: 94% accuracy, 0.89 balanced.',
    'Without calibration the model just answers “occupied” (0.52 balanced) — '
    'calibrating against the empty room is what makes CSI usable.',
    'Activity (5 classes) is far harder at ~0.50, though still ahead of RSSI-only '
    '(0.42) and raw CSI (0.38).',
], size=16)

# ---- Slide 4: presence confusion ------------------------------------------
s = slide(); title(s, 'Presence: confusion matrix')
image_fit(s, os.path.join(A, 'fig2_presence_confusion.png'), 0.6, 1.5, 6.2, 5.4)
bullets(s, 7.1, 2.2, 5.7, 4.5, [
    'Occupied is caught 96% of the time (1077 / 1127).',
    'Empty is confirmed 83% of the time (116 / 140).',
    'The residual error is calling a truly empty room “occupied” — the harder '
    'direction.',
    'Numbers pooled over all five held-out configurations.',
])

# ---- Slide 5: per-class recall --------------------------------------------
s = slide(); title(s, 'Where activity classification breaks down')
image_fit(s, os.path.join(A, 'fig3_perclass_recall.png'), 0.6, 1.5, 6.4, 5.4)
bullets(s, 7.2, 2.2, 5.6, 4.5, [
    'Empty, standing and walking are recognizable (0.53–0.88 recall).',
    'Sitting (0.28) is confused with standing; running (0.14) with walking.',
    'A single 2.4 GHz link cannot separate fine motion classes.',
    'This is the main argument for adding a 2nd and 3rd receiver.',
])

# ---- Slide 6: baseline sensitivity ----------------------------------------
s = slide(); title(s, 'Calibration detail: the empty baseline')
image_fit(s, os.path.join(A, 'fig4_baseline_sensitivity.png'), 0.6, 1.5, 6.4, 5.4)
bullets(s, 7.2, 2.2, 5.6, 4.5, [
    'The empty-room signal is stable within a session (>0.999 correlation).',
    'But a short one-shot empty capture underestimates normal variation and '
    'over-reports occupancy (empty recall 0.32).',
    'A longer, representative empty baseline restores it (0.83).',
    'Deployment rule: record a longer baseline at install, not one snapshot.',
])

# ---- Slide 7: status & next -----------------------------------------------
s = slide(); title(s, 'Status & next steps')
bullets(s, 0.7, 1.5, SW - 1.4, 5.5, [
    'Interim snapshot: one room, one subject, one session per configuration — '
    'no confidence intervals yet.',
    'Solid today: presence detection, robust across placements once calibrated.',
    'Not yet: fine activity classification with a single link.',
    'Next: more configurations (for confidence intervals), more sit/run data, then '
    'a 2nd and 3rd receiver for motion coverage, followed by other rooms and subjects.',
], size=19)

out = 'docs/Phase_A_Interim_Slides.pptx'
prs.save(out)
print('Saved', out, '-', len(prs.slides._sldIdLst), 'slides')
