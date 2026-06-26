#!/usr/bin/env python3
"""Build the weekly-progress PowerPoint (simple, workplace style)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
LIGHT = RGBColor(0xE8, 0xF0, 0xF8)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def title_bar(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = ACCENT
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.65), Inches(1.18), Inches(3.0), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background()


def bullets(tf, items):
    tf.word_wrap = True
    for i, it in enumerate(items):
        text, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(17 if lvl == 0 else 14)
        r.font.color.rgb = DARK if lvl == 0 else GRAY
        p.space_after = Pt(7)


def caption(slide, left, top, width, text):
    c = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    c.text_frame.word_wrap = True
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GRAY


def heading(slide, left, top, width, text, size=18):
    h = slide.shapes.add_textbox(left, top, width, Inches(0.45))
    p = h.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = ACCENT


def content(title, items, image=None, cap=None, img_w=Inches(5.8)):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title)
    if image:
        tb = s.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(6.2), Inches(5.6))
        bullets(tb.text_frame, items)
        pic = s.shapes.add_picture(image, Inches(7.0), Inches(1.7), width=img_w)
        pic.left = int(Inches(7.0) + (Inches(6.0) - pic.width) / 2)
        if cap:
            caption(s, Inches(7.0), Inches(1.7) + pic.height + Inches(0.08), Inches(6.0), cap)
    else:
        tb = s.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(12), Inches(5.6))
        bullets(tb.text_frame, items)
    return s


def image_slide(title, image, cap=None, width=Inches(11.5)):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title)
    pic = s.shapes.add_picture(image, Inches(1.0), Inches(1.7), width=width)
    pic.left = int((SW - pic.width) / 2)
    if cap:
        caption(s, Inches(0.8), Inches(1.7) + pic.height + Inches(0.15),
                SW - Inches(1.6), cap)
    return s


def stats_table(slide):
    data = [
        ['', 'Same setup', 'New placement'],
        ['Presence accuracy', '81%', '73%'],
        ['Occupied recall (catches a person)', '0.94', '0.86'],
        ['Empty recall', '0.61', '0.53'],
        ['Balanced accuracy', '0.77', '0.70'],
    ]
    tbl = slide.shapes.add_table(len(data), 3, Inches(1.0), Inches(1.7),
                                 Inches(11.3), Inches(2.6)).table
    tbl.columns[0].width = Inches(5.3)
    tbl.columns[1].width = Inches(3.0)
    tbl.columns[2].width = Inches(3.0)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(15)
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = (RGBColor(0xFF, 0xFF, 0xFF) if ri % 2
                                            else RGBColor(0xEE, 0xF3, 0xF9))


def pipeline(slide, steps, top=Inches(5.1)):
    n = len(steps)
    margin = Inches(0.55)
    gap = Inches(0.28)
    box_w = int((SW - 2 * margin - gap * (n - 1)) / n)
    bh = Inches(1.25)
    x = margin
    for i, st in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, bh)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = ACCENT; box.line.width = Pt(1.25)
        tf = box.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = st
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = DARK
        if i < n - 1:
            ar = slide.shapes.add_textbox(x + box_w, top, gap, bh)
            ar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            ap = ar.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
            arr = ap.add_run(); arr.text = "→"
            arr.font.size = Pt(20); arr.font.color.rgb = ACCENT
        x += box_w + gap


# ============================ TITLE (only) ============================
s = prs.slides.add_slide(BLANK)
box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2))
r = box.text_frame.paragraphs[0].add_run(); r.text = "Wi-Fi CSI Presence Sensing"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = ACCENT
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.85), Inches(4.5), Pt(3))
rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background()
sub = s.shapes.add_textbox(Inches(0.95), Inches(4.0), Inches(11.5), Inches(1.2))
for i, line in enumerate(["Weekly progress — what a single ESP32 can and cannot do",
                          "Dylan Dhawan  •  June 25, 2026"]):
    pp = sub.text_frame.paragraphs[0] if i == 0 else sub.text_frame.add_paragraph()
    rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(20 if i == 0 else 15)
    rr.font.color.rgb = DARK if i == 0 else GRAY

# ============== GOAL & WHAT IS CSI (combined, two columns) ==============
s = prs.slides.add_slide(BLANK)
title_bar(s, "Goal & What is CSI")
heading(s, Inches(0.65), Inches(1.5), Inches(6.0), "The goal")
ltb = s.shapes.add_textbox(Inches(0.65), Inches(2.05), Inches(6.0), Inches(5.0))
bullets(ltb.text_frame, [
    "Turn ordinary Wi-Fi into a presence / motion sensor — no cameras or wearables.",
    "Roadmap: presence → activity → people-counting → coarse localization.",
    "Master 1 ESP link first, then scale to 2 and 3.",
    "Build methods that can transfer to better radios later.",
])
heading(s, Inches(7.05), Inches(1.5), Inches(6.0), "What is CSI?")
rtb = s.shapes.add_textbox(Inches(7.05), Inches(2.05), Inches(5.9), Inches(5.0))
bullets(rtb.text_frame, [
    "Wi-Fi (20 MHz) splits into 64 subcarriers via OFDM (~52 usable).",
    "Each packet → the channel measured per subcarrier (amplitude + phase).",
    "That matrix is the Channel State Information (CSI).",
    "Far richer than RSSI (a single signal-strength number).",
    "The ESP32-C3 exposes raw CSI — a ~$5 chip becomes a sensor.",
])

# ============================ HOW WE GET INFO ============================
image_slide("How we read information from CSI", "slides_assets/spectrogram_compare.png",
            cap="Wi-Fi reflects off walls, furniture, and bodies (multipath). A static room gives stable CSI; "
                "a person or motion constantly changes it. We sample CSI ~100x/sec and analyze the change over time. "
                "Left: empty room (stable bands). Right: a moving person (rapid changes / Doppler ripples).")

# ============================ TELLING ACTIONS APART ============================
image_slide("Telling actions apart in the CSI", "slides_assets/spectrogram3.png",
            cap="Empty = stable. A still, standing person = a subtle change — the model relies on small, slow fluctuations "
                "rather than obvious motion. Moving = large, rapid, broadband changes. These differences are what the model learns to separate.")

# ============== SETUP & DATA COLLECTION (combined) ==============
content("Our Setup & Data Collection", [
    "2x Seeed XIAO ESP32-C3 (2.4 GHz, 20 MHz): one transmits, one receives — a single Wi-Fi link.",
    "Custom firmware: fixed channel, MAC-filtered, steady 100 Hz; 0 dropped packets to the laptop.",
    "A webcam records alongside the CSI, time-stamped on the same clock.",
    "Pose estimation (Ultralytics YOLO11-pose, COCO-pretrained) auto-labels each moment empty / still / moving.",
    "Result: thousands of time-aligned, labeled CSI windows with almost no manual effort.",
    "Tested across 3 link placements in one conference room.",
], image="photos/Picture of setup.jpg", cap="TX and RX boards on the wall; laptop logs CSI between them.")

# ============================ MODEL & TRAINING ============================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The model & how we train it")
tb = s.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(12.1), Inches(3.4))
bullets(tb.text_frame, [
    "CSI is sliced into 2-second windows (~200 samples x 52 subcarriers).",
    "Features per window: amplitude variability (motion), deviation from an “empty” baseline (presence), "
    "and low-frequency fluctuation energy 0.15-0.5 Hz (helps catch still occupants).",
    "Per-setup calibration: record a few seconds of empty at install, then subtract & scale-normalize — "
    "this is what lets one model work across different placements.",
    "Classifier: Random Forest — robust on modest data and interpretable.",
    "Validation: train on some setups, test on a completely unseen one (leave-one-out) — measures real-world generalization, "
    "not memorization of one room.",
])
pipeline(s, ["CSI stream\n@ 100 Hz", "2-second\nwindows",
             "Features:\nmotion, low-freq,\nempty-deviation",
             "Per-setup\ncalibration", "Random Forest\n→ empty / still / moving"])

# ============================ WHAT WORKS ============================
content("What 1 ESP is good at", [
    "Reliable presence detection — including a still, standing person (no movement required).",
    "~81% accuracy when tested within a setup it trained on.",
    "~73% on a brand-new link placement, after a quick empty-room calibration.",
    "Very high “someone is here” recall — rarely misses an occupant.",
    "Generalizes across placements once calibrated — a key step toward portability.",
], image="slides_assets/presence_bar.png")

# ============================ FALLS FLAT ============================
content("Where 1 ESP falls flat", [
    "Motion is only sensed near the link — detection drops to ~0% as a person moves away.",
    "Distinguishing motion far from the link is unreliable (the radio simply doesn’t “see” it).",
    "People-counting and precise location are not dependable with a single link.",
    "One link covers only a slice of the room — large blind spots remain.",
], image="slides_assets/coverage_bar.png")

# ============== WHY MORE ESPS + NEXT STEPS (combined) ==============
s = prs.slides.add_slide(BLANK)
title_bar(s, "Why more ESPs — and what's next")
heading(s, Inches(0.65), Inches(1.5), Inches(6.0), "Why we need more ESPs")
ltb = s.shapes.add_textbox(Inches(0.65), Inches(2.05), Inches(6.0), Inches(5.0))
bullets(ltb.text_frame, [
    "Each link covers a different region; together they fill each other’s blind spots.",
    "Fusing 2-3 links → whole-room motion & position coverage, people counting, robustness.",
    "Coverage data is the concrete evidence: one link physically cannot cover the room.",
    "More viewpoints also stabilize presence accuracy across the space.",
])
heading(s, Inches(7.05), Inches(1.5), Inches(6.0), "Next steps")
rtb = s.shapes.add_textbox(Inches(7.05), Inches(2.05), Inches(5.8), Inches(5.0))
bullets(rtb.text_frame, [
    "Add the 2nd ESP; directly measure the coverage / accuracy gain.",
    "More setups + a cleaner generalization test to firm up the ~73% cross-placement number.",
    "Begin fusing multiple links; explore people-counting and coarse localization.",
    "Short methodology write-up: capture → auto-labeling → calibration recipe.",
])

# ============================ APPENDIX DIVIDER ============================
s = prs.slides.add_slide(BLANK)
box = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.0))
r = box.text_frame.paragraphs[0].add_run(); r.text = "Appendix"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = ACCENT
sub = s.shapes.add_textbox(Inches(0.95), Inches(3.9), Inches(11.5), Inches(0.6))
rr = sub.text_frame.paragraphs[0].add_run(); rr.text = "Backup detail — pull up only if asked"
rr.font.size = Pt(16); rr.font.color.rgb = GRAY

# ---- A1: Model & training details ----
content("Appendix — Model & training details", [
    "Trained on 12 sessions (~35-40 min of CSI) → ~1,900 labeled 2-second windows.",
    "3 link placements, one conference room, one person; labels auto-generated by webcam + YOLO-pose.",
    "Features per window: amplitude variability (motion), deviation from an empty baseline (presence), "
    "low-frequency energy 0.15-0.5 Hz, and a per-setup calibration.",
    "Classifier: Random Forest (400 trees) — strong on modest tabular data, interpretable, no scaling needed.",
    "Validation: leave-one-out; hardest test trains on 2 placements and tests on a brand-new 3rd.",
    "To try later: gradient boosting (XGBoost/LightGBM); CNN/LSTM with much more data.",
])

# ---- A2: Accuracy & stats ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "Appendix — Accuracy & stats")
stats_table(s)
note = s.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.3), Inches(2.2))
bullets(note.text_frame, [
    "3-class (empty / still / moving): ~61-67%.",
    "Motion (moving vs. still): weak (~0.2-0.4 recall); ~0% far from the link.",
    "A naive random-split test would read ~99% (memorizes the room) — we report held-out numbers.",
    "Model leans cautious: rarely misses a person, more likely to false-alarm an empty room.",
])

# ---- A3: OFDM ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "Appendix — OFDM, the basis of CSI")
tb = s.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(12.1), Inches(1.7))
bullets(tb.text_frame, [
    "Wi-Fi splits its 20 MHz channel into 64 subcarriers (~52 usable), sent in parallel — that is OFDM.",
    "“Orthogonal” = the tones are spaced so they don’t interfere (each peak sits on its neighbours’ zeros).",
    "The receiver measures the channel at each subcarrier → those ~52 values ARE the CSI (vs. RSSI’s one number).",
])
pic = s.shapes.add_picture("slides_assets/ofdm_diagram.png", Inches(1.3), Inches(3.3), width=Inches(10.7))
pic.left = int((SW - pic.width) / 2)

# ---- A4: External datasets & models ----
content("Appendix — External datasets & models", [
    "Public CSI datasets exist: Widar3.0, SignFi, UT-HAR, FallDeFi, WiAR.",
    "Caveat: mostly different hardware (e.g. Intel 5300: 30 subcarriers, 3 antennas) and rooms — "
    "CSI is hardware- and environment-specific.",
    "We showed our own model breaks just by moving the link, so external data won’t directly raise our ESP32 accuracy.",
    "Useful for: feature/architecture ideas, benchmarking, and later transfer-learning — not plug-and-play.",
    "No good drop-in pretrained CSI model for ESP32; the value online is in methods, not weights.",
])

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "CSI Weekly Progress.pptx"
prs.save(out)
print(f"Saved {out} with", len(prs.slides._sldIdLst), "slides")
